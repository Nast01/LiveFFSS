"""Génère le support de démonstration LiveFFSS (docs/demo/livefss-demo.pptx).

Régénère le fichier depuis zéro : toute retouche faite dans PowerPoint sera perdue.
Les captures d'écran sont lues dans docs/demo/screenshots/ et rognées à la volée
(barre d'état Android en haut, barre de geste en bas).

    python docs/demo/build_deck.py
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
CROPPED = ROOT / ".cropped"
OUT = ROOT / "livefss-demo.pptx"

# Charte reprise de lib/app/core/theme/app_colors.dart
BLUE = RGBColor(0x21, 0x96, 0xF3)
BLUE_DARK = RGBColor(0x19, 0x76, 0xD2)
ORANGE = RGBColor(0xFB, 0x8C, 0x00)
GREEN = RGBColor(0x43, 0xA0, 0x47)
RED = RGBColor(0xE5, 0x39, 0x35)
GREY = RGBColor(0x9E, 0x9E, 0x9E)
INK = RGBColor(0x21, 0x21, 0x21)
INK_SOFT = RGBColor(0x61, 0x61, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF7, 0xF9, 0xFB)
TODO_BG = RGBColor(0xFF, 0xF4, 0xCC)
TODO_FG = RGBColor(0x8A, 0x62, 0x00)

FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)

# Statuts affichés en pastille sur les slides d'écran.
LIVRE = ("Livré", GREEN)
PARTIEL = ("Partiel", ORANGE)
COQUILLE = ("Coquille", GREY)
ABSENT = ("Non implémenté", RED)


# --------------------------------------------------------------------------- #
# primitives


def crop(name):
    """Rogne la barre d'état et la barre de geste, renvoie le chemin du PNG."""
    CROPPED.mkdir(exist_ok=True)
    src, dst = SHOTS / f"{name}.png", CROPPED / f"{name}.png"
    if not src.exists():
        return None
    if not dst.exists() or dst.stat().st_mtime < src.stat().st_mtime:
        im = Image.open(src)
        im.crop((0, 100, im.width, im.height - 40)).save(dst)
    return dst


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = liste de (texte, taille, gras, couleur, espace_avant_pt)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, bold, color, space) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space)
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def badge(slide, x, y, label, color, size=11):
    w = Inches(0.28 + 0.085 * len(label))
    s = rect(slide, x, y, w, Inches(0.32), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return s


def footer(slide, n):
    text(
        slide,
        Inches(0.5),
        H - Inches(0.5),
        Inches(6),
        Inches(0.3),
        [("LiveFFSS · Démonstration", 9, False, GREY, 0)],
    )
    text(
        slide,
        W - Inches(1.2),
        H - Inches(0.5),
        Inches(0.7),
        Inches(0.3),
        [(str(n), 9, False, GREY, 0)],
        align=PP_ALIGN.RIGHT,
    )


def todo_banner(slide, x, y, w, message):
    b = rect(slide, x, y, w, Inches(0.45), TODO_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"  [À COMPLÉTER]  {message}"
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = TODO_FG


def header(slide, title, kicker=None, status=None):
    rect(slide, 0, 0, W, Inches(1.0), BLUE)
    runs = []
    if kicker:
        runs.append((kicker.upper(), 10, True, RGBColor(0xCC, 0xE5, 0xFA), 0))
    runs.append((title, 24, True, WHITE, 0))
    text(slide, Inches(0.6), Inches(0.12), Inches(9.5), Inches(0.8), runs)
    if status:
        badge(slide, W - Inches(2.4), Inches(0.34), status[0], status[1])


# --------------------------------------------------------------------------- #
# gabarits de slides


def slide_title(prs, n):
    s = blank(prs)
    rect(s, 0, 0, W, H, BLUE)
    rect(s, 0, H - Inches(1.4), W, Inches(1.4), BLUE_DARK)
    text(
        s,
        Inches(1.0),
        Inches(2.2),
        Inches(11),
        Inches(2.2),
        [
            ("LiveFFSS", 66, True, WHITE, 0),
            (
                "Suivi et organisation des compétitions de sauvetage sportif",
                22,
                False,
                RGBColor(0xE3, 0xF2, 0xFD),
                14,
            ),
        ],
    )
    text(
        s,
        Inches(1.0),
        H - Inches(1.15),
        Inches(11),
        Inches(0.9),
        [("Démonstration  ·  [DATE]  ·  [PRÉSENTÉ PAR]", 15, False, WHITE, 0)],
    )
    return s


def slide_section(prs, n, number, title, subtitle):
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    rect(s, 0, Inches(2.6), Inches(0.22), Inches(1.9), ORANGE)
    text(
        s,
        Inches(0.8),
        Inches(2.5),
        Inches(11),
        Inches(2.2),
        [
            (number, 15, True, ORANGE, 0),
            (title, 44, True, INK, 6),
            (subtitle, 17, False, INK_SOFT, 10),
        ],
    )
    footer(s, n)
    return s


def slide_screen(prs, n, kicker, title, purpose, bullets, shot, status, note=None):
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    header(s, title, kicker, status)

    img = crop(shot) if shot else None
    img_h = Inches(5.55)
    if img:
        with Image.open(img) as im:
            ratio = im.width / im.height
        img_w = Emu(int(img_h * ratio))
        pic = s.shapes.add_picture(str(img), Inches(0.75), Inches(1.35), height=img_h)
        pic.line.color.rgb = RGBColor(0xDD, 0xE3, 0xE8)
        pic.line.width = Pt(0.75)
    else:
        img_w = Inches(2.65)
        ph = rect(
            s,
            Inches(0.75),
            Inches(1.35),
            img_w,
            img_h,
            RGBColor(0xEC, 0xF1, 0xF5),
            line=RGBColor(0xC5, 0xCF, 0xD6),
        )
        tf = ph.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "Capture à insérer"
        r.font.name = FONT
        r.font.size = Pt(13)
        r.font.color.rgb = GREY

    left = Inches(0.75) + img_w + Inches(0.7)
    width = W - left - Inches(0.75)

    text(s, left, Inches(1.45), width, Inches(1.0), [(purpose, 15, False, INK, 0)])

    runs = []
    for b in bullets:
        warn = b.startswith("!")
        runs.append(
            (
                ("⚠  " if warn else "•  ") + (b[1:].strip() if warn else b),
                13,
                False,
                ORANGE if warn else INK_SOFT,
                10,
            )
        )
    text(s, left, Inches(2.75), width, Inches(3.6), runs)

    if note:
        todo_banner(s, left, H - Inches(1.35), width, note)

    footer(s, n)
    return s


def slide_bullets(prs, n, kicker, title, intro, blocks, note=None):
    """blocks = liste de (titre, [lignes]) rendue en deux colonnes."""
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    header(s, title, kicker)
    y = Inches(1.35)
    if intro:
        text(s, Inches(0.75), y, W - Inches(1.5), Inches(0.6), [(intro, 15, False, INK, 0)])
        y = Inches(2.05)

    cols = 2 if len(blocks) > 2 else 1
    col_w = (W - Inches(1.5) - Inches(0.6) * (cols - 1)) / cols
    per_col = -(-len(blocks) // cols)

    for i, (btitle, lines) in enumerate(blocks):
        c, r_ = divmod(i, per_col)
        x = Inches(0.75) + c * (col_w + Inches(0.6))
        by = y + r_ * Inches(1.62)
        rect(s, x, by + Inches(0.06), Inches(0.06), Inches(0.28), ORANGE)
        runs = [(btitle, 15, True, INK, 0)]
        for line in lines:
            warn = line.startswith("!")
            runs.append(
                (
                    ("⚠  " if warn else "·  ") + (line[1:].strip() if warn else line),
                    12,
                    False,
                    ORANGE if warn else INK_SOFT,
                    7,
                )
            )
        text(s, x + Inches(0.22), by, col_w - Inches(0.22), Inches(1.5), runs)

    if note:
        todo_banner(s, Inches(0.75), H - Inches(1.35), W - Inches(1.5), note)
    footer(s, n)
    return s


def slide_map(prs, n):
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    header(s, "Carte de l'application", "vue d'ensemble")

    def node(x, y, w, h, label, sub, fill, fg=WHITE):
        b = rect(s, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = b.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = fg
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = sub
            r2.font.name = FONT
            r2.font.size = Pt(10)
            r2.font.color.rgb = fg
        return b

    def arrow(x1, y1, x2, y2):
        c = s.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = ligne droite
        c.line.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)
        c.line.width = Pt(1.75)

    col = [Inches(0.75), Inches(4.05), Inches(7.35), Inches(10.65)]
    cw = Inches(2.55)

    node(col[0], Inches(1.6), cw, Inches(1.0), "Accueil", "recherche · filtres", BLUE)
    node(col[0], Inches(2.9), cw, Inches(0.85), "Favoris", "", BLUE)
    node(col[0], Inches(4.05), cw, Inches(0.85), "Connexion", "profil", GREY)

    node(
        col[1],
        Inches(1.6),
        cw,
        Inches(2.25),
        "Détail compétition",
        "Évènements · Programme\nClubs · Points",
        BLUE_DARK,
    )
    node(col[1], Inches(4.35), cw, Inches(0.95), "Bracelets RFID", "écriture NFC", ORANGE)
    node(col[1], Inches(5.6), cw, Inches(0.95), "Programme", "structure · horaires", ORANGE)

    node(
        col[2],
        Inches(1.6),
        cw,
        Inches(1.8),
        "Détail épreuve",
        "Engagés · Séries\nRésumé",
        BLUE_DARK,
    )
    node(col[2], Inches(5.6), cw, Inches(0.95), "Éditeur de structure", "tours · bracket", ORANGE)

    node(col[3], Inches(1.6), cw, Inches(0.95), "Tirage des séries", "", ORANGE)
    node(col[3], Inches(2.85), cw, Inches(0.95), "Course", "résultats à venir", GREY)

    arrow(col[0] + cw, Inches(2.1), col[1], Inches(2.4))
    arrow(col[0] + cw, Inches(3.3), col[1], Inches(2.8))
    arrow(col[1] + cw, Inches(2.4), col[2], Inches(2.2))
    arrow(col[1] + Inches(1.3), Inches(3.85), col[1] + Inches(1.3), Inches(4.35))
    arrow(col[1] + Inches(1.3), Inches(5.3), col[1] + Inches(1.3), Inches(5.6))
    arrow(col[1] + cw, Inches(6.05), col[2], Inches(6.05))
    arrow(col[2] + cw, Inches(2.1), col[3], Inches(2.1))
    arrow(col[2] + cw, Inches(2.8), col[3], Inches(3.3))

    text(
        s,
        Inches(0.75),
        H - Inches(1.05),
        W - Inches(1.5),
        Inches(0.5),
        [
            (
                "Bleu : consultation, ouverte à tous.   Orange : organisation, "
                "réservée aux organisateurs.   Gris : écran incomplet.",
                12,
                False,
                INK_SOFT,
                0,
            )
        ],
    )
    footer(s, n)
    return s


def slide_table(prs, n, title, kicker, headers, rows, widths):
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    header(s, title, kicker)

    left, top = Inches(0.7), Inches(1.3)
    total_w = W - Inches(1.4)
    row_h = Inches(0.325)
    shape = s.shapes.add_table(
        len(rows) + 1, len(headers), left, top, total_w, row_h * (len(rows) + 1)
    )
    table = shape.table
    for i, frac in enumerate(widths):
        table.columns[i].width = Emu(int(total_w * frac))
    for row in table.rows:
        row.height = row_h

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE

    palette = {"Livré": GREEN, "Partiel": ORANGE, "Coquille": GREY, "Non implémenté": RED}
    for r_, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r_, c)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_ % 2 else RGBColor(0xEE, 0xF3, 0xF7)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.bold = val in palette
            run.font.color.rgb = palette.get(val, INK_SOFT)
    footer(s, n)
    return s


# --------------------------------------------------------------------------- #
# contenu


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    n = 0

    def step():
        nonlocal n
        n += 1
        return n

    slide_title(prs, step())

    slide_bullets(
        prs,
        step(),
        "contexte",
        "Pourquoi LiveFFSS",
        "",
        [
            ("Aujourd'hui", ["[Comment se déroule l'organisation actuellement]"]),
            ("Les points de friction", ["[Ce qui coince : papier, tableurs, attente…]"]),
            ("Ce qu'on veut obtenir", ["[Objectif de l'application]"]),
            ("Qui s'en sert", ["[Organisateurs, jury, athlètes, public]"]),
        ],
        note="cette page est une trame — remplace les quatre blocs par ton propre contexte.",
    )

    slide_bullets(
        prs,
        step(),
        "en bref",
        "LiveFFSS en une page",
        "Une application mobile Flutter branchée sur l'API de la FFSS (ffss.fr), "
        "pour suivre une compétition de sauvetage sportif et pour l'organiser.",
        [
            (
                "Deux usages",
                [
                    "Consulter : compétitions, épreuves, engagés, clubs",
                    "Organiser : programme, séries, marshalling, bracelets",
                ],
            ),
            (
                "Deux disciplines",
                ["Côtier (plage) — 16 places par course", "Eau plate (bassin) — 8 places"],
            ),
            (
                "Sans compte / avec compte",
                [
                    "Tout le volet consultation est ouvert",
                    "La connexion FFSS débloque les actions d'organisation",
                ],
            ),
            (
                "Chiffres",
                [
                    "12 écrans, 287 fichiers Dart",
                    "66 fichiers de tests unitaires",
                    "Interface bilingue FR / EN",
                ],
            ),
        ],
    )

    slide_map(prs, step())

    slide_bullets(
        prs,
        step(),
        "convention",
        "Comment lire les pastilles",
        "Chaque écran de la démo porte une pastille d'état, en haut à droite.",
        [
            ("Livré", ["L'écran fait ce qu'on attend de lui, avec de vraies données."]),
            (
                "Partiel",
                [
                    "Fonctionne, mais avec une limite assumée — le plus souvent :",
                    "les données restent sur l'appareil, faute d'endpoint FFSS.",
                ],
            ),
            (
                "Coquille",
                [
                    "L'écran existe et se navigue, son contenu reste à écrire.",
                ],
            ),
            (
                "Non implémenté",
                [
                    "L'interface est prête, la source de données est bouchonnée.",
                ],
            ),
        ],
    )

    # ---------------------------------------------------------------- parcours
    slide_section(
        prs, step(), "01", "Parcours spectateur", "Ce que voit quelqu'un qui ouvre l'app sans compte."
    )

    slide_screen(
        prs,
        step(),
        "accueil",
        "Trouver une compétition",
        "Le point d'entrée : toutes les compétitions FFSS, la plus proche en tête.",
        [
            "Trois vues : Vus récemment · Cette semaine · Tous",
            "« Cette semaine » est filtrée côté serveur, pas dans l'app",
            "Filtre de discipline : Tous / Eau plate / Côtier",
            "Étoile pour suivre, pastille Ouvert / Fermé pour les engagements",
            "Tirer vers le bas pour rafraîchir",
        ],
        "03-accueil-tous",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "accueil",
        "Recherche",
        "Recherche instantanée sur le nom et le lieu, appliquée à la vue courante.",
        [
            "Se combine avec le filtre de discipline",
            "Chaque carte porte la date, le lieu et le logo du club organisateur",
            "À défaut de logo, un bloc date lisible de loin",
        ],
        "04-accueil-recherche",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "favoris",
        "Suivre ses compétitions",
        "Deuxième onglet de la barre inférieure : les compétitions mises en favori.",
        [
            "Aucun compte nécessaire — la liste vit dans le stockage sécurisé du téléphone",
            "Même carte que l'accueil, retrait d'un tap",
            "L'app garde aussi les 20 dernières compétitions consultées",
        ],
        "23-favoris",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "compte",
        "Connexion",
        "Identifiant et mot de passe FFSS. Le jeton est conservé en stockage sécurisé.",
        [
            "La consultation ne demande jamais de compte",
            "La connexion débloque le programme, le tirage et les bracelets",
            "Expiration de session gérée globalement : déconnexion et retour ici",
            "! Cet écran est encore sur le thème Material par défaut — à aligner sur la charte",
        ],
        "02-connexion",
        PARTIEL,
    )

    slide_screen(
        prs,
        step(),
        "compte",
        "Profil",
        "Identité du licencié : nom, club, rôle, et déconnexion.",
        [
            "Alimenté par l'endpoint FFSS du profil connecté",
            "Le profil est mis en cache pour un démarrage hors ligne",
            "! Capture manquante : l'appareil de démo n'était pas connecté",
        ],
        None,
        LIVRE,
        note="ouvre l'app connecté, capture l'écran Profil et dépose-le ici.",
    )

    slide_screen(
        prs,
        step(),
        "réglages",
        "Bilingue",
        "Français et anglais, le choix est persisté entre deux lancements.",
        [
            "Toutes les chaînes d'interface passent par le catalogue de traductions",
            "Aucun texte en dur dans les écrans",
        ],
        "24-langue",
        LIVRE,
    )

    # ------------------------------------------------------------- compétition
    slide_section(
        prs,
        step(),
        "02",
        "La compétition",
        "Quatre onglets : ce qui se court, quand, avec qui, et pour combien de points.",
    )

    slide_screen(
        prs,
        step(),
        "compétition · onglet 1",
        "Évènements",
        "L'en-tête porte les dates, le lieu, le club organisateur et les accès organisateur.",
        [
            "Une épreuve = discipline × sexe × catégorie",
            "Filtre Côtier / Eau plate",
            "Icône calendrier : construction du programme",
            "Icône NFC : écriture des bracelets",
            "Étoile : mise en favori depuis la fiche",
        ],
        "05-competition-epreuves",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "compétition · onglet 3",
        "Clubs",
        "Tous les clubs engagés, avec leur effectif d'athlètes et d'arbitres.",
        [
            "Recherche transverse : club, athlète ou arbitre",
            "Logo du club, puis bonnet, puis initiale — jamais d'icône générique",
            "Les clubs étrangers et les invités sont résolus eux aussi",
        ],
        "07-competition-clubs",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "compétition · onglet 3",
        "Clubs — le détail",
        "Déplier un club donne la liste nominative de ses athlètes et arbitres.",
        [
            "Utile au pointage d'arrivée d'une délégation",
            "Même source que la liste des engagés d'une épreuve",
        ],
        "09-clubs-detail",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "compétition · onglet 4",
        "Points",
        "Trois classements prévus : par club, individuel, et relais.",
        [
            "L'écran, la navigation et l'état vide sont en place",
            "! Les endpoints de classement FFSS ne sont pas documentés",
            "! La source de données renvoie des listes vides — rien à afficher",
            "Le jour où le backend arrive, seule la source change",
        ],
        "08-competition-points",
        ABSENT,
    )

    slide_screen(
        prs,
        step(),
        "compétition · onglet 2",
        "Programme (lecture)",
        "La vue publique du déroulé : par jour, par site, heure par heure.",
        [
            "Une ligne par course, avec son créneau et sa durée",
            "Liseré bleu pour les séries, vert pour les finales",
            "! Alimenté par le programme construit dans l'app, pas par le serveur",
        ],
        "06b-competition-programme-defini",
        PARTIEL,
    )

    # ------------------------------------------------------------------ épreuve
    slide_section(
        prs,
        step(),
        "03",
        "L'épreuve et le marshalling",
        "Là où se joue la journée : qui est là, dans quelle série, et à quelle heure.",
    )

    slide_screen(
        prs,
        step(),
        "épreuve · onglet 1",
        "Engagés",
        "Les athlètes inscrits sur l'épreuve, avec club et année de naissance.",
        [
            "Tri par nom, par club ou par dossard",
            "Chaque athlète porte un état de présence",
            "Rafraîchissement automatique toutes les 10 secondes",
            "Le bonnet du club sert de vignette",
        ],
        "10-epreuve-engages",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "épreuve · onglet 1",
        "Pointer les présents",
        "Le cœur du marshalling : valider qui se présente au départ.",
        [
            "Un tap bascule « En attente » en « Présent »",
            "Ou scan du bracelet NFC, sans toucher la liste",
            "Le pointage alimente directement le tirage des séries",
            "Conservé entre deux rafraîchissements et deux ouvertures de l'app",
            "! Stocké sur l'appareil : deux téléphones qui pointent la même épreuve "
            "ne voient pas le travail de l'autre",
        ],
        "10b-engages-presence",
        PARTIEL,
    )

    slide_screen(
        prs,
        step(),
        "bracelets",
        "Écrire les bracelets",
        "Encodage NFC des bracelets, en amont de la compétition.",
        [
            "Tous les athlètes de la compétition, recherche par nom",
            "Chaque ligne affiche la charge utile qui sera écrite",
            "Doublons entre clubs dédupliqués automatiquement",
            "Accessible depuis l'en-tête de la compétition",
        ],
        "21-rfid-writer",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "bracelets",
        "Le format du bracelet",
        "licence | sexe | type | département | club ; NOM — relu tel quel au marshalling.",
        [
            "Écriture au contact, une pression sur l'athlète sélectionné",
            "! L'émulateur n'a pas de puce NFC, d'où le message d'erreur sur la capture",
            "Sur téléphone équipé, l'écriture et le scan fonctionnent",
            "Un bracelet mal écrit peut être réécrit sans limite",
        ],
        "22-rfid-ecriture",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "épreuve · onglet 2",
        "Séries",
        "La structure de l'épreuve : les tours, leurs courses, leurs places.",
        [
            "Séries, quarts, demi-finales, finale selon ce qui est défini",
            "« Générer les séries » lance le tirage sur les présents",
            "Ouvrir une série donne accès à sa course",
            "Reste vide tant qu'aucune structure n'est définie",
        ],
        "11-epreuve-series",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "épreuve · onglet 3",
        "Résumé",
        "Vue de synthèse d'une épreuve une fois courue.",
        [
            "! Dépend des résultats, qui ne sont pas disponibles",
            "L'écran et son état vide existent, le contenu reste à écrire",
        ],
        "12-epreuve-resume",
        COQUILLE,
    )

    slide_screen(
        prs,
        step(),
        "course",
        "Une course",
        "Ouverte depuis une série : épreuve, catégorie, niveau et numéro de course.",
        [
            "C'est ici que se feront la saisie et l'affichage des résultats",
            "! Les endpoints de résultats FFSS ne sont pas documentés",
            "Le contexte est déjà transmis correctement — il ne manque que les données",
        ],
        "20-course",
        COQUILLE,
    )

    # ---------------------------------------------------------------- programme
    slide_section(
        prs,
        step(),
        "04",
        "Construire le programme",
        "De la liste des engagés à une grille horaire, en trois étapes.",
    )

    slide_screen(
        prs,
        step(),
        "programme · structure",
        "L'inventaire",
        "Une ligne par épreuve × catégorie, avec son nombre d'engagés.",
        [
            "« Aucun déroulement FFSS » = le serveur ne connaît pas encore le format",
            "Deux actions en masse : générer toutes les structures, "
            "créer les déroulements manquants côté FFSS",
            "C'est la seule écriture réelle de l'app vers le serveur, hors connexion",
        ],
        "13-programme-structure-vide",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "programme · structure",
        "Génération automatique",
        "L'app propose un format à partir du nombre d'engagés et des places par course.",
        [
            "16 places par course en côtier, 8 en eau plate",
            "Si la FFSS connaît déjà un déroulement, il sert de base",
            "Chaque ligne reste modifiable individuellement",
        ],
        "14-programme-structure",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "programme · structure",
        "Éditeur de structure",
        "Le format d'une épreuve, tour par tour.",
        [
            "Niveaux : séries, quarts, demi-finales, finale",
            "Par niveau : nombre de courses, places par course, qualifiés par course",
            "« Proposer une structure » recalcule tout d'un coup",
            "Supprimer un niveau le supprime aussi côté FFSS",
        ],
        "15-structure-editor",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "programme · structure",
        "Vue bracket",
        "La même structure, lue comme un tableau de qualification.",
        [
            "Contrôle visuel avant de figer le format",
            "Bascule liste / bracket depuis l'en-tête",
        ],
        "16-structure-bracket",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "programme · horaires",
        "Les sites",
        "Un site est un lieu de course : plage, zone de sable, bassin.",
        [
            "Plusieurs sites peuvent tourner en parallèle sur une même journée",
            "Le type conditionne les épreuves qu'on y planifie",
        ],
        "17b-programme-sites",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "programme · horaires",
        "La grille horaire",
        "Par jour et par site : la file des courses, avec heure de départ et durée.",
        [
            "Heure de départ de la journée réglable",
            "Réordonnancement par glisser-déposer, durée ajustable course par course",
            "Les horaires se recalculent en cascade à chaque changement",
            "Les courses non planifiées restent listées en bas de l'écran",
        ],
        "17-programme-horaires",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "tirage",
        "Le tirage des séries",
        "Répartit les athlètes présents dans les courses du tour choisi.",
        [
            "Le compteur « présents sur engagés » rend la décision explicite",
            "Répartition équilibrée entre les courses, ordre aléatoire",
            "Échange de deux athlètes de couloir, relance possible du tirage",
            "Le niveau se choisit en haut : séries, puis finale",
        ],
        "18-tirage-series",
        LIVRE,
    )

    slide_screen(
        prs,
        step(),
        "tirage",
        "Séries enregistrées",
        "Une fois sauvegardé, le tirage est visible depuis l'épreuve.",
        [
            "Chaque série affiche son effectif réel",
            "Le programme et la vue publique reflètent immédiatement le tirage",
            "! La FFSS n'expose aucun endpoint d'écriture des séries : "
            "le tirage vit sur l'appareil qui l'a fait",
        ],
        "19-series-apres-tirage",
        PARTIEL,
    )

    # ---------------------------------------------------------------- technique
    slide_section(
        prs, step(), "05", "Sous le capot", "Trois pages pour situer l'architecture et ses limites."
    )

    slide_bullets(
        prs,
        step(),
        "technique",
        "Architecture",
        "Flutter / Dart, trois couches, une règle : chaque domaine possède les siennes.",
        [
            (
                "Controller → Repository → DataSource",
                [
                    "Le controller porte l'état de l'écran, jamais d'appel réseau",
                    "Le repository orchestre (pagination, composition)",
                    "La source de données parle HTTP et rien d'autre",
                ],
            ),
            (
                "Du JSON au domaine",
                [
                    "DTO calqués sur le JSON FFSS, champs en français",
                    "Un mapper les traduit en modèles typés, en anglais",
                    "Modèles immuables, sérialisation générée",
                ],
            ),
            (
                "Sept domaines",
                [
                    "auth · club · compétition · épreuve",
                    "meeting · résultat · classement",
                    "Un seul point d'injection au démarrage",
                ],
            ),
            (
                "Conventions tenues",
                [
                    "Pas de contexte graphique ni de traduction dans les controllers",
                    "Injection par constructeur uniquement",
                    "Analyse statique stricte, pas de conversions implicites",
                ],
            ),
        ],
    )

    slide_bullets(
        prs,
        step(),
        "technique",
        "L'API FFSS : ce qu'on peut, ce qu'on ne peut pas",
        "L'API est externe et figée. Elle décide de ce que l'app sait faire.",
        [
            (
                "Ce qu'on lit",
                [
                    "Compétitions, épreuves, engagements",
                    "Clubs, athlètes, arbitres",
                    "Déroulements (formats d'épreuve)",
                    "Profil du licencié connecté",
                ],
            ),
            (
                "Ce qu'on écrit",
                [
                    "Authentification",
                    "Création et suppression de déroulement",
                    "…et c'est tout.",
                ],
            ),
            (
                "Ce qui n'existe pas",
                [
                    "! Résultats et temps",
                    "! Classements de compétition",
                    "! Séries et compositions",
                    "! Présence au marshalling",
                ],
            ),
            (
                "La conséquence",
                [
                    "Tout ce qui touche au déroulé du jour J reste sur l'appareil.",
                    "C'est le verrou principal, et il n'est pas de notre côté.",
                ],
            ),
        ],
    )

    slide_bullets(
        prs,
        step(),
        "technique",
        "Ce qui vit sur le téléphone",
        "Faute d'endpoints, plusieurs briques sont persistées localement, en stockage sécurisé.",
        [
            (
                "Ce qui est stocké",
                [
                    "Jeton d'authentification et profil",
                    "Favoris et compétitions récemment consultées",
                    "Présence au marshalling, par épreuve",
                    "Programme complet : structures, sites, horaires, séries",
                ],
            ),
            (
                "Les limites à dire en démo",
                [
                    "! Rien n'est partagé entre deux appareils",
                    "! Une réinstallation efface le programme construit",
                    "! Pas de reprise si le téléphone du chrono tombe",
                ],
            ),
            (
                "Tests",
                [
                    "66 fichiers de tests unitaires",
                    "Mappers, repositories, controllers, client HTTP",
                    "Pas de tests d'interface : vérification manuelle",
                ],
            ),
            (
                "Ce que ça coûterait de lever",
                [
                    "Un endpoint d'écriture programme/séries côté FFSS",
                    "Le reste de l'app est déjà découpé pour l'accueillir",
                ],
            ),
        ],
    )

    # ------------------------------------------------------------------- reste
    slide_section(
        prs, step(), "06", "Ce qu'il reste à faire", "L'état écran par écran, puis les chantiers."
    )

    slide_table(
        prs,
        step(),
        "État écran par écran",
        "récapitulatif",
        ["Écran", "État", "Ce qui bloque"],
        [
            ["Accueil · Favoris · Recherche", "Livré", "—"],
            ["Connexion", "Partiel", "Thème par défaut, à aligner sur la charte"],
            ["Profil", "Livré", "—"],
            ["Compétition · Évènements", "Livré", "—"],
            ["Compétition · Clubs", "Livré", "—"],
            ["Compétition · Points", "Non implémenté", "Endpoints de classement non documentés"],
            ["Compétition · Programme", "Partiel", "Lit le programme local, pas le serveur"],
            ["Épreuve · Engagés", "Livré", "—"],
            ["Marshalling (présence, scan)", "Partiel", "Présence locale, non partagée"],
            ["Épreuve · Séries", "Livré", "—"],
            ["Épreuve · Résumé", "Coquille", "Dépend des résultats"],
            ["Course", "Coquille", "Endpoints de résultats non documentés"],
            ["Programme · Structure + éditeur", "Livré", "—"],
            ["Programme · Sites + horaires", "Livré", "Persistance locale uniquement"],
            ["Tirage des séries", "Partiel", "Aucune écriture serveur possible"],
            ["Bracelets RFID", "Livré", "Nécessite un téléphone NFC"],
        ],
        [0.34, 0.16, 0.50],
    )

    slide_bullets(
        prs,
        step(),
        "reste à faire",
        "Chantiers fonctionnels",
        "Par ordre d'impact sur une compétition réelle.",
        [
            (
                "1 · Résultats live",
                [
                    "Saisie des temps de nage et des classements de plage",
                    "Consultation des résultats d'une course, abandons",
                    "! Quatre appels à câbler dès que la FFSS documente ses endpoints",
                ],
            ),
            (
                "2 · Partage entre appareils",
                [
                    "Programme, séries et présence sur un serveur commun",
                    "Sans ça, un seul téléphone pilote la journée",
                ],
            ),
            (
                "3 · Classements / Points",
                [
                    "Clubs, individuel, relais",
                    "L'interface est prête, il manque la source",
                ],
            ),
            (
                "4 · Écrans à remplir",
                [
                    "Course : affichage et saisie des résultats",
                    "Résumé d'épreuve",
                ],
            ),
        ],
    )

    slide_bullets(
        prs,
        step(),
        "reste à faire",
        "Chantiers techniques",
        "Rien de bloquant, mais à traiter avant d'ouvrir le code à d'autres contributeurs.",
        [
            (
                "Cohérence visuelle",
                [
                    "L'écran de connexion est encore sur le thème par défaut",
                    "Quelques dialogues utilisent la palette Material et non la charte",
                ],
            ),
            (
                "Dette identifiée",
                [
                    "L'écran de créneau hérité fait 634 lignes — à découper",
                    "Il pilote encore ses dialogues depuis le controller",
                ],
            ),
            (
                "Typage",
                [
                    "La discipline est reconnue par comparaison de chaînes",
                    "À remplacer par une énumération typée",
                ],
            ),
            (
                "Ménage",
                [
                    "Trois routes déclarées mais jamais branchées",
                    "Trois controllers historiques hors du canal de messages unifié",
                ],
            ),
        ],
    )

    slide_bullets(
        prs,
        step(),
        "suite",
        "Prochaines étapes",
        "",
        [
            ("Court terme", ["[Ce que tu veux livrer dans les prochaines semaines]"]),
            ("Ce qu'il faut de la FFSS", ["[Endpoints à obtenir, interlocuteur, échéance]"]),
            ("Test en conditions réelles", ["[Quelle compétition, quelle date]"]),
            ("Moyens", ["[Temps, matériel NFC, personnes]"]),
        ],
        note="trame à remplir : la roadmap dépend de ce que la FFSS peut ouvrir côté API.",
    )

    s = blank(prs)
    rect(s, 0, 0, W, H, BLUE)
    text(
        s,
        Inches(1.0),
        Inches(3.0),
        Inches(11),
        Inches(1.5),
        [
            ("Questions", 54, True, WHITE, 0),
            ("Merci.", 20, False, RGBColor(0xE3, 0xF2, 0xFD), 12),
        ],
    )

    prs.save(OUT)
    print(f"{OUT} : {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
